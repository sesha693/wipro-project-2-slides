from dataclasses import dataclass
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import config as cfg

# Local fallbacks for config values to avoid AttributeError if user config is missing entries
BG_COLOR = getattr(cfg, "BACKGROUND_COLOR", "#FFFFFF")
HEADER_FONT = getattr(cfg, "HEADER_FONT", "Calibri")
HEADER_FONT_SIZE = getattr(cfg, "HEADER_FONT_SIZE", 28)
HEADER_ACCENT_COLOR = getattr(cfg, "HEADER_ACCENT_COLOR", "#1F3A57")
COLOR_TABLE_HEADER = getattr(cfg, "COLOR_TABLE_HEADER", "#F3F6F8")
TABLE_BORDER_COLOR = getattr(cfg, "TABLE_BORDER_COLOR", "#D9D9D9")
PANEL_BG_COLOR = getattr(cfg, "PANEL_BG_COLOR", "#FFFFFF")
CARD_BG = getattr(cfg, "CARD_BG", "#FFFFFF")
CARD_BORDER_WIDTH = getattr(cfg, "CARD_BORDER_WIDTH", 1.5)
CARD_BORDER_COLOR = getattr(cfg, "CARD_BORDER_COLOR", "#D9E2EA")
SLIDE_WIDTH = getattr(cfg, "SLIDE_WIDTH_INCHES", 13.3333)
SLIDE_HEIGHT = getattr(cfg, "SLIDE_HEIGHT_INCHES", 7.5)

# Populate missing attributes on cfg so other code referencing cfg.<NAME> won't fail
_defaults = {
    "BACKGROUND_COLOR": BG_COLOR,
    "HEADER_FONT": HEADER_FONT,
    "HEADER_FONT_SIZE": HEADER_FONT_SIZE,
    "HEADER_ACCENT_COLOR": HEADER_ACCENT_COLOR,
    "COLOR_TABLE_HEADER": COLOR_TABLE_HEADER,
    "TABLE_BORDER_COLOR": TABLE_BORDER_COLOR,
    "PANEL_BG_COLOR": PANEL_BG_COLOR,
    "CARD_BG": CARD_BG,
    "CARD_BORDER_WIDTH": CARD_BORDER_WIDTH,
    "CARD_BORDER_COLOR": CARD_BORDER_COLOR,
    "SLIDE_WIDTH_INCHES": SLIDE_WIDTH,
    "SLIDE_HEIGHT_INCHES": SLIDE_HEIGHT,
}
for _k, _v in _defaults.items():
    if not hasattr(cfg, _k):
        setattr(cfg, _k, _v)

DISPLAY_HEADER_MAP = {
    "adh": "ADH",
    "account": "Account",
    "q1_27": "Q1'27",
    "q2_27_target": "Q2'27 Target",
    "q2_27": "Q2'27",
    "actuals": "Actuals",
    "target": "Target",
    "people_cost": "People Cost",
    "non_people_cost": "Non-People Cost",
    "prev_week_total": "Prev Week Total",
    "curr_week_total": "Curr Week Total",
    "prev_week_people_cost": "Prev Week People",
    "prev_week_non_people_cost": "Prev Week Non-People",
    "curr_week_people_cost": "Curr Week People",
    "curr_week_non_people_cost": "Curr Week Non-People",
    "q1_gm": "Q1 GM",
    "q2_gm_target": "Q2 GM Target",
    "gm_forecast": "GM Forecast",
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


@dataclass
class Card:
    label: str
    value: str
    color_hex: str


@dataclass
class SlideData:
    title: str
    table_columns: list[str]
    table_values: list[str]
    cards: list[Card]
    chart_png: bytes
    notes: Optional[str] = None
    table_title: str = "Service Line View"
    chart_title: str = "Service Line Movement Matrix"
    callouts: Optional[list[str]] = None


def add_background(slide):
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(cfg.SLIDE_WIDTH_INCHES),
        Inches(cfg.SLIDE_HEIGHT_INCHES),
    ).fill.solid()
    bg = slide.shapes[-1]
    bg.fill.fore_color.rgb = _hex_to_rgb(cfg.BACKGROUND_COLOR)
    bg.line.fill.background()


def add_header(slide, text: str, left: Inches, top: Inches, width: Inches, height: Inches):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = cfg.HEADER_FONT
    p.font.size = Pt(cfg.HEADER_FONT_SIZE)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(cfg.HEADER_ACCENT_COLOR)
    p.alignment = PP_ALIGN.LEFT
    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top + height - Inches(0.12),
        width * 0.18,
        Inches(0.09),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _hex_to_rgb(cfg.HEADER_ACCENT_COLOR)
    accent_bar.line.fill.background()


def add_single_row_table(slide, columns: list[str], values: list[str], left, top, width, height):
    padding = Inches(0.05)
    container = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left - padding,
        top - padding,
        width + padding * 2,
        height + padding * 2,
    )
    container.fill.solid()
    container.fill.fore_color.rgb = RGBColor(250, 250, 250)
    container.line.color.rgb = _hex_to_rgb(cfg.TABLE_BORDER_COLOR)
    container.line.width = Pt(1)

    table_shape = slide.shapes.add_table(rows=2, cols=len(columns), left=left, top=top, width=width, height=height)
    table = table_shape.table
    header_color = _hex_to_rgb(cfg.COLOR_TABLE_HEADER)
    for idx, header in enumerate(columns):
        cell = table.cell(0, idx)
        cell.text = DISPLAY_HEADER_MAP.get(header, header).replace("_", " ").title()
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(9)
        para.alignment = PP_ALIGN.CENTER
        cell.text_frame.margin_left = Pt(3)
        cell.text_frame.margin_right = Pt(3)
        cell.text_frame.margin_top = Pt(3)
        cell.text_frame.margin_bottom = Pt(3)

    for idx, value in enumerate(values):
        cell = table.cell(1, idx)
        cell.text = value
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11)
        para.alignment = PP_ALIGN.CENTER
        cell.text_frame.margin_left = Pt(3)
        cell.text_frame.margin_right = Pt(3)
        cell.text_frame.margin_top = Pt(3)
        cell.text_frame.margin_bottom = Pt(3)


def add_data_cards(slide, cards: list[Card], left, top, width, height, cols=4):
    card_width = width / cols
    card_height = height
    for index, card in enumerate(cards):
        card_left = left + card_width * index
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            card_left,
            top,
            card_width - Inches(0.1),
            card_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(cfg.CARD_BG)
        shape.line.color.rgb = _hex_to_rgb(card.color_hex)
        shape.line.width = Pt(cfg.CARD_BORDER_WIDTH)
        text_frame = shape.text_frame
        text_frame.clear()
        label_para = text_frame.paragraphs[0]
        label_para.text = card.label.upper()
        label_para.font.size = Pt(9)
        label_para.font.bold = True
        label_para.font.color.rgb = _hex_to_rgb(cfg.HEADER_ACCENT_COLOR)
        label_para.alignment = PP_ALIGN.LEFT
        label_para.space_after = Pt(4)
        value_para = text_frame.add_paragraph()
        value_para.text = card.value
        value_para.font.size = Pt(18)
        value_para.font.bold = True
        value_para.font.color.rgb = _hex_to_rgb(card.color_hex)
        value_para.alignment = PP_ALIGN.LEFT


def add_panel_container(slide, left, top, width, height):
    container = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    container.fill.solid()
    container.fill.fore_color.rgb = _hex_to_rgb(cfg.PANEL_BG_COLOR)
    container.line.color.rgb = _hex_to_rgb(CARD_BORDER_COLOR)
    container.line.width = Pt(0.75)
    return container


def add_panel_title(slide, text: str, left, top, width, height):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = cfg.HEADER_FONT
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb("#2C3E50")
    p.alignment = PP_ALIGN.LEFT


def add_panel_table(slide, title: str, columns: list[str], values: list[str], left, top, width, height):
    add_panel_container(slide, left, top, width, height)
    title_height = Inches(0.28)
    add_panel_title(slide, title, left + Inches(0.15), top + Inches(0.14), width - Inches(0.3), title_height)
    table_top = top + title_height + Inches(0.14)
    table_height = height - title_height - Inches(0.18)
    add_single_row_table(slide, columns, values, left + Inches(0.1), table_top, width - Inches(0.2), table_height)


def add_panel_chart(slide, title: str, chart_png: bytes, left, top, width, height):
    add_panel_container(slide, left, top, width, height)
    title_height = Inches(0.24)
    add_panel_title(slide, title, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), title_height)
    add_chart_image(slide, chart_png, left + Inches(0.1), top + title_height + Inches(0.12), width - Inches(0.2), height - title_height - Inches(0.16))


def add_cards_panel(slide, cards: list[Card], left, top, width, height):
    add_panel_container(slide, left, top, width, height)
    add_data_cards(slide, cards, left + Inches(0.07), top + Inches(0.08), width - Inches(0.14), height - Inches(0.16), cols=len(cards))


def create_pptx(slide_data_list: list[SlideData], filename: str, template_path: Optional[str] = None) -> bytes:
    prs = Presentation(template_path) if template_path else Presentation()
    prs.slide_width = Inches(cfg.SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(cfg.SLIDE_HEIGHT_INCHES)
    for slide_data in slide_data_list:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        add_background(slide)
        add_header(slide, slide_data.title, Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6))
        add_cards_panel(slide, slide_data.cards, Inches(0.3), Inches(0.95), Inches(12.6), Inches(1.05))
        add_panel_table(slide, slide_data.table_title, slide_data.table_columns, slide_data.table_values, Inches(0.3), Inches(2.15), Inches(7.0), Inches(3.6))
        add_panel_chart(slide, slide_data.chart_title, slide_data.chart_png, Inches(7.4), Inches(2.15), Inches(5.5), Inches(3.6))
        if slide_data.callouts:
            callout_top = Inches(5.85)
            add_panel_container(slide, Inches(0.3), callout_top, Inches(12.6), Inches(1.0))
            text_box = slide.shapes.add_textbox(Inches(0.4), callout_top + Inches(0.15), Inches(12.4), Inches(0.7))
            frame = text_box.text_frame
            frame.clear()
            p = frame.paragraphs[0]
            p.text = "Key Insights"
            p.font.name = cfg.HEADER_FONT
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(cfg.HEADER_ACCENT_COLOR)
            p.space_after = Pt(4)
            for callout in slide_data.callouts:
                p = frame.add_paragraph()
                p.text = f"• {callout}"
                p.font.size = Pt(11)
                p.font.name = cfg.HEADER_FONT
                p.space_after = Pt(2)
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.notes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
